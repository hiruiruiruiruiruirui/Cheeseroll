"""Document parsing service using LangChain document loaders.

Supports: .pptx, .docx, .pdf
Falls back to lightweight parsers (python-pptx, python-docx, pymupdf) when
the full Unstructured library is not installed.
"""

import os
from typing import Any


async def parse_document(file_path: str, file_type: str) -> dict[str, Any]:
    """Parse a document file and return structured content.

    Args:
        file_path: Path to the document file on disk.
        file_type: One of 'pptx', 'docx', 'pdf'.

    Returns:
        Dict with keys: title (str), sections (list of {heading, level, content, tables})
    """
    parsers = {
        "pptx": _parse_pptx,
        "docx": _parse_docx,
        "pdf": _parse_pdf,
    }

    parser = parsers.get(file_type)
    if not parser:
        raise ValueError(f"Unsupported file type: {file_type}")

    return await parser(file_path)


def _extract_title(filename: str) -> str:
    """Extract a readable title from a filename."""
    name = os.path.splitext(os.path.basename(filename))[0]
    # Remove common prefixes/suffixes
    return name.strip()


async def _parse_pptx(file_path: str) -> dict[str, Any]:
    """Parse PowerPoint files.

    Try UnstructuredPowerPointLoader first, fall back to python-pptx.
    """
    try:
        from langchain_community.document_loaders import UnstructuredPowerPointLoader

        loader = UnstructuredPowerPointLoader(file_path, mode="elements")
        docs = loader.load()

        title = _extract_title(file_path)
        sections: list[dict[str, Any]] = []
        current_section: dict[str, Any] | None = None

        for doc in docs:
            text = doc.page_content.strip()
            if not text:
                continue

            metadata = doc.metadata
            category = metadata.get("category", "")

            if category == "Title":
                if current_section:
                    sections.append(current_section)
                current_section = {
                    "heading": text,
                    "level": 2,
                    "content": "",
                    "tables": [],
                }
            elif current_section is not None:
                if category == "Table":
                    current_section.setdefault("tables", []).append(text)
                else:
                    current_section["content"] += text + "\n"
            else:
                current_section = {
                    "heading": title,
                    "level": 1,
                    "content": text + "\n",
                    "tables": [],
                }

        if current_section:
            sections.append(current_section)

        return {"title": title, "sections": sections}

    except ImportError:
        # Fallback to python-pptx
        return await _parse_pptx_lightweight(file_path)


async def _parse_pptx_lightweight(file_path: str) -> dict[str, Any]:
    """Lightweight PPTX parser using python-pptx only."""
    from pptx import Presentation

    prs = Presentation(file_path)
    title = _extract_title(file_path)
    sections: list[dict[str, Any]] = []
    all_text: list[str] = []

    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
            if shape.has_table:
                table = shape.table
                rows: list[str] = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append("| " + " | ".join(cells) + " |")
                if rows:
                    if len(rows) > 1:
                        rows.insert(1, "|" + "|".join(["---"] * len(table.columns)) + "|")
                    slide_texts.append("\n".join(rows))

        if slide_texts:
            heading = slide_texts[0] if slide_texts else f"幻灯片 {slide_idx}"
            content = "\n".join(slide_texts[1:]) if len(slide_texts) > 1 else ""
            sections.append({
                "heading": heading,
                "level": 2,
                "content": content,
                "tables": [],
            })
            all_text.extend(slide_texts)

    if not sections:
        sections.append({"heading": title, "level": 1, "content": "\n".join(all_text), "tables": []})

    return {"title": title, "sections": sections}


async def _parse_docx(file_path: str) -> dict[str, Any]:
    """Parse Word documents.

    Try UnstructuredWordDocumentLoader first, fall back to python-docx.
    """
    try:
        from langchain_community.document_loaders import UnstructuredWordDocumentLoader

        loader = UnstructuredWordDocumentLoader(file_path, mode="elements")
        docs = loader.load()

        title = _extract_title(file_path)
        sections: list[dict[str, Any]] = []
        current_section: dict[str, Any] | None = None

        for doc in docs:
            text = doc.page_content.strip()
            if not text:
                continue

            category = doc.metadata.get("category", "")
            if category == "Title":
                if current_section:
                    sections.append(current_section)
                current_section = {
                    "heading": text,
                    "level": 2,
                    "content": "",
                    "tables": [],
                }
            elif current_section is not None:
                current_section["content"] += text + "\n"
            else:
                current_section = {
                    "heading": title,
                    "level": 1,
                    "content": text + "\n",
                    "tables": [],
                }

        if current_section:
            sections.append(current_section)

        return {"title": title, "sections": sections}

    except ImportError:
        return await _parse_docx_lightweight(file_path)


async def _parse_docx_lightweight(file_path: str) -> dict[str, Any]:
    """Lightweight DOCX parser using python-docx only."""
    from docx import Document

    doc = Document(file_path)
    title = _extract_title(file_path)
    sections: list[dict[str, Any]] = []
    current_heading = title
    current_content: list[str] = []
    current_level = 1

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        if para.style.name.startswith("Heading"):
            if current_content:
                sections.append({
                    "heading": current_heading,
                    "level": current_level,
                    "content": "\n".join(current_content),
                    "tables": [],
                })
                current_content = []

            current_heading = text
            try:
                current_level = int(para.style.name.split()[-1])
            except ValueError:
                current_level = 2
        else:
            current_content.append(text)

    # Don't forget the last section
    if current_content or current_heading != title:
        sections.append({
            "heading": current_heading,
            "level": current_level,
            "content": "\n".join(current_content),
            "tables": [],
        })

    if not sections:
        sections.append({"heading": title, "level": 1, "content": "", "tables": []})

    return {"title": title, "sections": sections}


async def _parse_pdf(file_path: str) -> dict[str, Any]:
    """Parse PDF files.

    Try UnstructuredPDFLoader first, fall back to pymupdf (fitz).
    """
    try:
        from langchain_community.document_loaders import UnstructuredPDFLoader

        loader = UnstructuredPDFLoader(file_path, mode="elements")
        docs = loader.load()

        title = _extract_title(file_path)
        sections: list[dict[str, Any]] = []
        current_section: dict[str, Any] | None = None

        for doc in docs:
            text = doc.page_content.strip()
            if not text:
                continue

            category = doc.metadata.get("category", "")
            if category == "Title":
                if current_section:
                    sections.append(current_section)
                current_section = {
                    "heading": text,
                    "level": 2,
                    "content": "",
                    "tables": [],
                }
            elif current_section is not None:
                current_section["content"] += text + "\n"
            else:
                current_section = {
                    "heading": title,
                    "level": 1,
                    "content": text + "\n",
                    "tables": [],
                }

        if current_section:
            sections.append(current_section)

        return {"title": title, "sections": sections}

    except ImportError:
        return await _parse_pdf_lightweight(file_path)


async def _parse_pdf_lightweight(file_path: str) -> dict[str, Any]:
    """Lightweight PDF parser using pymupdf (fitz)."""
    import fitz  # pymupdf

    doc = fitz.open(file_path)
    title = _extract_title(file_path)
    all_pages_text: list[str] = []

    for page in doc:
        text = page.get_text("text")
        if text.strip():
            all_pages_text.append(text.strip())

    doc.close()

    full_text = "\n\n".join(all_pages_text)

    # Split by double newlines into rough sections
    sections: list[dict[str, Any]] = []
    paragraphs = full_text.split("\n\n")
    current_content: list[str] = []
    current_heading = title

    for para in paragraphs:
        para = para.strip()
        if not para:
            continue
        # Heuristic: short lines without punctuation might be headings
        if len(para) < 80 and "\n" not in para and not para.endswith((".", "。", "！", "？")):
            if current_content:
                sections.append({
                    "heading": current_heading,
                    "level": 2,
                    "content": "\n\n".join(current_content),
                    "tables": [],
                })
                current_content = []
            current_heading = para
        else:
            current_content.append(para)

    if current_content:
        sections.append({
            "heading": current_heading,
            "level": 2,
            "content": "\n\n".join(current_content),
            "tables": [],
        })

    if not sections:
        sections.append({"heading": title, "level": 1, "content": full_text, "tables": []})

    return {"title": title, "sections": sections}
